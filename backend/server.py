from fastapi import FastAPI, APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse, Response
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import json
import asyncio
import subprocess
import base64
import glob
from pathlib import Path
from pydantic import BaseModel, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone

from pptx import Presentation as PptxPresentation
from docx import Document as DocxDocument
import fitz  # PyMuPDF
from openai import AsyncOpenAI
from PIL import Image
import pytesseract
import io

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY', '')
UPLOAD_DIR = Path("/tmp/pptx_uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

app = FastAPI()
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pptx", ".docx", ".pdf"}


# ════════════════════════════════════════
# Models
# ════════════════════════════════════════
class JobInfo(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    filename: str
    file_type: str
    total_segments: int
    status: str

class TranslateRequest(BaseModel):
    target_language: str
    tone: str = "formal"

class ProgressResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    status: str
    progress: float
    total_segments: int
    translated_segments: int
    error_message: Optional[str] = None


# ════════════════════════════════════════
# PPTX extraction & rebuild
# ════════════════════════════════════════
def extract_pptx_segments(path: str) -> list:
    prs = PptxPresentation(path)
    segments = []
    idx = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(r.text for r in para.runs)
                    if full_text.strip():
                        segments.append({
                            "idx": idx,
                            "slide_num": slide_idx + 1,
                            "original": full_text,
                            "translated": None,
                            "translatable": not full_text.strip().startswith("http"),
                        })
                        idx += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            full_text = "".join(r.text for r in para.runs)
                            if full_text.strip():
                                segments.append({
                                    "idx": idx,
                                    "slide_num": slide_idx + 1,
                                    "original": full_text,
                                    "translated": None,
                                    "translatable": True,
                                })
                                idx += 1
    return segments


def rebuild_pptx(original_path: str, output_path: str, translations: dict):
    prs = PptxPresentation(original_path)
    idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(r.text for r in para.runs)
                    if full_text.strip():
                        translated = translations.get(idx)
                        if translated:
                            runs = para.runs
                            if runs:
                                runs[0].text = translated
                                for r in runs[1:]:
                                    r.text = ""
                        idx += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            full_text = "".join(r.text for r in para.runs)
                            if full_text.strip():
                                translated = translations.get(idx)
                                if translated:
                                    runs = para.runs
                                    if runs:
                                        runs[0].text = translated
                                        for r in runs[1:]:
                                            r.text = ""
                                idx += 1
    prs.save(output_path)


# ════════════════════════════════════════
# DOCX extraction & rebuild
# ════════════════════════════════════════
def extract_docx_segments(path: str) -> list:
    doc = DocxDocument(path)
    segments = []
    idx = 0
    page_num = 1
    line_count = 0

    # Extract paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            segments.append({
                "idx": idx,
                "slide_num": page_num,  # approximate page number
                "original": text,
                "translated": None,
                "translatable": not text.startswith("http"),
            })
            idx += 1
            line_count += 1
            # Rough page estimation (every ~35 lines)
            if line_count > 35:
                page_num += 1
                line_count = 0

    # Extract table cells
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    # Avoid duplicates from merged cells
                    if not segments or segments[-1]["original"] != text:
                        segments.append({
                            "idx": idx,
                            "slide_num": page_num,
                            "original": text,
                            "translated": None,
                            "translatable": True,
                        })
                        idx += 1
    return segments


def rebuild_docx(original_path: str, output_path: str, translations: dict):
    doc = DocxDocument(original_path)
    idx = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            translated = translations.get(idx)
            if translated and para.runs:
                para.runs[0].text = translated
                for r in para.runs[1:]:
                    r.text = ""
            idx += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    translated = translations.get(idx)
                    if translated:
                        for p in cell.paragraphs:
                            if p.runs:
                                p.runs[0].text = translated
                                for r in p.runs[1:]:
                                    r.text = ""
                                translated = ""  # Only set in first paragraph of cell
                    idx += 1
    doc.save(output_path)


# ════════════════════════════════════════
# PDF extraction & rebuild
# ════════════════════════════════════════
def extract_pdf_segments(path: str) -> tuple:
    """Extract text segments from PDF. Returns (segments, is_ocr).
    Falls back to OCR if no text layer is found.
    """
    doc = fitz.open(path)
    segments = []
    idx = 0
    for page_num, page in enumerate(doc):
        blocks = page.get_text("blocks")
        for block in blocks:
            # block: (x0, y0, x1, y1, text, block_no, block_type)
            if block[6] == 0:  # text block
                text = block[4].strip()
                if text and not text.startswith("http"):
                    segments.append({
                        "idx": idx,
                        "slide_num": page_num + 1,
                        "original": text,
                        "translated": None,
                        "translatable": True,
                    })
                    idx += 1
    doc.close()

    # If no text found, fall back to OCR
    if not segments:
        logger.info(f"No text layer found in PDF, falling back to OCR: {path}")
        segments = extract_pdf_segments_ocr(path)
        return segments, True

    return segments, False


def extract_pdf_segments_ocr(path: str) -> list:
    """Extract text segments from image-based PDF using OCR (Tesseract).
    Uses line-level bounding boxes with background color sampling for precise overlay rebuild.
    Returns segments with line-level detail for rebuilding.
    """
    doc = fitz.open(path)
    segments = []
    idx = 0
    render_dpi = 300

    for page_num, page in enumerate(doc):
        # Render page as high-res image for OCR
        pix = page.get_pixmap(dpi=render_dpi)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        img_width, img_height = img.size

        # OCR with bounding box data
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)

        # Group words into lines
        lines = {}
        for i in range(len(data['text'])):
            text = data['text'][i].strip()
            if text and data['conf'][i] > 30:
                block_num = data['block_num'][i]
                line_num = data['line_num'][i]
                key = (block_num, line_num)
                if key not in lines:
                    lines[key] = {
                        'words': [],
                        'left': data['left'][i],
                        'top': data['top'][i],
                        'right': data['left'][i] + data['width'][i],
                        'bottom': data['top'][i] + data['height'][i],
                        'height': data['height'][i],
                    }
                lines[key]['words'].append(text)
                lines[key]['left'] = min(lines[key]['left'], data['left'][i])
                lines[key]['top'] = min(lines[key]['top'], data['top'][i])
                lines[key]['right'] = max(lines[key]['right'], data['left'][i] + data['width'][i])
                lines[key]['bottom'] = max(lines[key]['bottom'], data['top'][i] + data['height'][i])
                lines[key]['height'] = max(lines[key]['height'], data['height'][i])

        # Helper: sample background color around a bounding box
        def sample_bg_color(left, top, right, bottom):
            samples = []
            offsets = [(-10, 0), (10, 0), (0, -5), (0, 5), (-10, -5), (10, 5)]
            for dx, dy in offsets:
                for sx, sy in [(left + dx, (top + bottom) // 2 + dy),
                               (right + dx, (top + bottom) // 2 + dy),
                               ((left + right) // 2, top + dy),
                               ((left + right) // 2, bottom + dy)]:
                    sx = max(0, min(int(sx), img_width - 1))
                    sy = max(0, min(int(sy), img_height - 1))
                    pixel = img.getpixel((sx, sy))
                    samples.append(pixel[:3])
            if not samples:
                return (255, 255, 255)
            # Use median to avoid outliers
            r = sorted([c[0] for c in samples])[len(samples) // 2]
            g = sorted([c[1] for c in samples])[len(samples) // 2]
            b = sorted([c[2] for c in samples])[len(samples) // 2]
            return (r, g, b)

        # Helper: detect text color (sample pixels within the text area)
        def sample_text_color(left, top, right, bottom):
            samples = []
            # Sample center pixels of the text area
            cx = (left + right) // 2
            cy = (top + bottom) // 2
            for dx in range(-20, 21, 10):
                for dy in range(-5, 6, 5):
                    sx = max(0, min(cx + dx, img_width - 1))
                    sy = max(0, min(cy + dy, img_height - 1))
                    pixel = img.getpixel((sx, sy))
                    samples.append(pixel[:3])
            if not samples:
                return (0, 0, 0)
            # Find darkest pixels (likely text)
            samples.sort(key=lambda c: sum(c))
            darkest = samples[:len(samples) // 3 + 1]
            r = sum(c[0] for c in darkest) // len(darkest)
            g = sum(c[1] for c in darkest) // len(darkest)
            b = sum(c[2] for c in darkest) // len(darkest)
            return (r, g, b)

        scale = 72.0 / render_dpi
        page_rect = page.rect

        # Build line-level segments for each block
        # Group lines by block for translation context
        block_lines = {}
        for (block_num, line_num), line_data in sorted(lines.items()):
            if block_num not in block_lines:
                block_lines[block_num] = []
            line_text = ' '.join(line_data['words'])
            if not line_text.strip() or len(line_text.strip()) < 2:
                continue

            # Get background and text colors
            bg_color = sample_bg_color(line_data['left'], line_data['top'],
                                        line_data['right'], line_data['bottom'])
            text_color = sample_text_color(line_data['left'], line_data['top'],
                                           line_data['right'], line_data['bottom'])

            # Convert pixel coordinates to PDF points
            x0 = line_data['left'] * scale
            y0 = line_data['top'] * scale
            x1 = line_data['right'] * scale
            y1 = line_data['bottom'] * scale
            line_height = line_data['height'] * scale

            # Clamp to page
            x0 = max(0, min(x0, page_rect.width))
            y0 = max(0, min(y0, page_rect.height))
            x1 = max(0, min(x1, page_rect.width))
            y1 = max(0, min(y1, page_rect.height))

            # Filter out OCR noise lines
            import re
            clean_text = line_text.strip()
            alpha_chars = sum(1 for c in clean_text if c.isalpha())
            total_chars = len(clean_text)
            if total_chars > 0 and alpha_chars / total_chars < 0.5:
                continue  # Skip lines with less than 50% alphabetic chars

            # Skip lines with mostly short garbage words (OCR noise)
            words = clean_text.split()
            if words:
                avg_word_len = sum(len(w) for w in words) / len(words)
                short_words = sum(1 for w in words if len(w) <= 2)
                # Noise patterns: short garbled words, pipe chars, excessive parens
                has_pipe = '|' in clean_text
                if len(words) >= 3 and avg_word_len < 3.5 and (short_words / len(words) > 0.3 or has_pipe):
                    continue  # OCR noise

            block_lines[block_num].append({
                'text': clean_text,
                'bbox': [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)],
                'line_height': round(line_height, 2),
                'bg_color': list(bg_color),
                'text_color': list(text_color),
            })

        # Create segments: one per block (for translation), but store line details
        for block_num in sorted(block_lines.keys()):
            blines = block_lines[block_num]
            if not blines:
                continue

            full_text = ' '.join(l['text'] for l in blines)
            if not full_text.strip() or full_text.startswith("http"):
                continue

            # Clean OCR artifacts from text
            # Remove common OCR noise prefixes/suffixes
            clean_full = full_text.strip()
            # Remove leading special chars like €), a , etc.
            clean_full = re.sub(r'^[€$£¥@#&*|<>\[\](){}\d]+[)\s.]+', '', clean_full).strip()
            # Remove trailing noise
            clean_full = re.sub(r'[\[\](){}<>|]+$', '', clean_full).strip()

            if not clean_full or len(clean_full) < 3:
                continue

            # Overall block bbox
            x0 = min(l['bbox'][0] for l in blines)
            y0 = min(l['bbox'][1] for l in blines)
            x1 = max(l['bbox'][2] for l in blines)
            y1 = max(l['bbox'][3] for l in blines)

            segments.append({
                "idx": idx,
                "slide_num": page_num + 1,
                "original": clean_full,
                "translated": None,
                "translatable": True,
                "ocr_bbox": [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)],
                "ocr_lines": blines,
            })
            idx += 1

    doc.close()
    logger.info(f"OCR extracted {len(segments)} text segments with line-level detail")
    return segments


def rebuild_pdf(original_path: str, output_path: str, translations: dict, is_ocr: bool = False):
    """Rebuild PDF by overlaying translated text on original blocks."""
    doc = fitz.open(original_path)

    # Use Liberation Sans for Unicode support (Croatian, etc.)
    FONT_PATH = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
    FONT_BOLD_PATH = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
    FONT_NAME = "libsans"
    FONT_BOLD_NAME = "libsansbold"

    if is_ocr:
        # For OCR-based (image) PDFs: overlay translated text directly on the original page
        # using line-level precision with background color matching.
        # translations dict: {idx: {"text": ..., "bbox": ..., "page": ..., "original": ..., "lines": [...]}}

        for idx_key, tdata in translations.items():
            if not isinstance(tdata, dict):
                continue
            page_idx = tdata.get("page", 0)
            if page_idx >= len(doc):
                continue
            page = doc[page_idx]
            translated = tdata.get("text", "")
            ocr_lines = tdata.get("lines", [])

            if not translated or not ocr_lines:
                continue

            # Split translated text into roughly the same number of lines as original
            original_lines_text = [l['text'] for l in ocr_lines]
            num_lines = len(ocr_lines)

            # Smart line splitting: try to distribute translated text across available lines
            translated_words = translated.split()
            if num_lines == 1:
                translated_lines = [translated]
            else:
                # Distribute words proportionally based on original line lengths
                orig_lengths = [len(lt) for lt in original_lines_text]
                total_orig_len = sum(orig_lengths) or 1
                translated_lines = []
                word_idx = 0
                for i, orig_len in enumerate(orig_lengths):
                    proportion = orig_len / total_orig_len
                    num_words = max(1, round(len(translated_words) * proportion))
                    if i == num_lines - 1:
                        # Last line gets remaining words
                        line_words = translated_words[word_idx:]
                    else:
                        line_words = translated_words[word_idx:word_idx + num_words]
                    translated_lines.append(' '.join(line_words))
                    word_idx += num_words

            # Overlay each translated line at the corresponding original line position
            for i, line_info in enumerate(ocr_lines):
                if i >= len(translated_lines) or not translated_lines[i].strip():
                    continue

                bbox = line_info['bbox']
                bg_color = line_info.get('bg_color', [255, 255, 255])
                text_color = line_info.get('text_color', [0, 0, 0])
                line_height = line_info.get('line_height', bbox[3] - bbox[1])

                # Create rect with small padding
                pad_x = 4
                pad_y = 2
                rect = fitz.Rect(
                    bbox[0] - pad_x,
                    bbox[1] - pad_y,
                    bbox[2] + pad_x,
                    bbox[3] + pad_y,
                )

                # Cover original text with background-colored rectangle
                fill_color = (bg_color[0] / 255.0, bg_color[1] / 255.0, bg_color[2] / 255.0)
                page.draw_rect(rect, color=fill_color, fill=fill_color)

                # Calculate font size based on line height
                font_size = min(line_height * 0.75, 60)
                font_size = max(font_size, 6)

                # Determine text color
                tc = (text_color[0] / 255.0, text_color[1] / 255.0, text_color[2] / 255.0)
                # If background is dark, use white text
                bg_brightness = (bg_color[0] + bg_color[1] + bg_color[2]) / 3
                if bg_brightness < 100:
                    tc = (1.0, 1.0, 1.0)
                elif bg_brightness > 200:
                    tc = (0.05, 0.05, 0.05)

                # Use bold font for larger text (likely headings)
                use_bold = font_size > 20
                fn = FONT_BOLD_NAME if use_bold else FONT_NAME
                ff = FONT_BOLD_PATH if use_bold else FONT_PATH

                # Insert translated text
                text_rect = fitz.Rect(bbox[0], bbox[1], bbox[2] + pad_x * 2, bbox[3] + pad_y)
                page.insert_textbox(
                    text_rect,
                    translated_lines[i],
                    fontsize=font_size,
                    fontname=fn,
                    fontfile=ff,
                    color=tc,
                    align=1,  # center align
                )

    else:
        # Standard text-layer PDF rebuild
        idx = 0
        for page in doc:
            blocks = page.get_text("blocks")
            for block in blocks:
                if block[6] == 0:
                    text = block[4].strip()
                    if text and not text.startswith("http"):
                        translated = translations.get(idx)
                        if translated:
                            rect = fitz.Rect(block[0], block[1], block[2], block[3])
                            page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))
                            rect_height = rect.height
                            font_size = min(rect_height * 0.8, 11)
                            font_size = max(font_size, 6)
                            page.insert_textbox(
                                rect,
                                translated,
                                fontsize=font_size,
                                fontname=FONT_NAME,
                                fontfile=FONT_PATH,
                                color=(0, 0, 0),
                                align=0,
                            )
                        idx += 1

    doc.save(output_path)
    doc.close()


# ════════════════════════════════════════
# Visual preview generation
# ════════════════════════════════════════
def generate_slide_images(file_path: str, file_type: str, job_id: str, suffix: str = "original") -> list:
    """Convert document to images for visual preview.
    Returns list of image file paths.
    """
    images_dir = UPLOAD_DIR / f"{job_id}_{suffix}_images"
    images_dir.mkdir(exist_ok=True)

    pdf_path = None

    if file_type == "pdf":
        pdf_path = file_path
    elif file_type in ("pptx", "docx"):
        # Convert to PDF using LibreOffice
        pdf_output_dir = UPLOAD_DIR / f"{job_id}_pdf"
        pdf_output_dir.mkdir(exist_ok=True)
        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", str(pdf_output_dir), file_path
                ],
                capture_output=True, text=True, timeout=120,
                env={**os.environ, "HOME": "/tmp"}
            )
            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                return []
            # Find the generated PDF
            base_name = Path(file_path).stem
            pdf_candidates = list(pdf_output_dir.glob(f"{base_name}*.pdf"))
            if not pdf_candidates:
                logger.error("No PDF generated by LibreOffice")
                return []
            pdf_path = str(pdf_candidates[0])
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
            return []
        except Exception as e:
            logger.error(f"LibreOffice conversion error: {e}")
            return []

    if not pdf_path:
        return []

    # Convert PDF to images using pdftoppm
    try:
        output_prefix = str(images_dir / "slide")
        subprocess.run(
            [
                "pdftoppm", "-png", "-r", "150",
                "-l", "10",  # Limit to first 10 pages for preview
                pdf_path, output_prefix
            ],
            capture_output=True, timeout=60
        )
        image_files = sorted(glob.glob(f"{output_prefix}*.png"))
        return image_files
    except Exception as e:
        logger.error(f"PDF to image conversion error: {e}")
        return []


# ════════════════════════════════════════
# Translation logic
# ════════════════════════════════════════
async def translate_batch(texts: list, target_language: str, tone: str) -> list:
    openai_client = AsyncOpenAI(api_key=OPENAI_API_KEY)
    tone_map = {
        "formal": "formal and professional",
        "academic": "academic and scholarly",
        "general": "general and natural-sounding",
    }
    tone_desc = tone_map.get(tone, "general and natural-sounding")

    prompt = f"""Translate the following text segments to {target_language}.
Tone: {tone_desc}.

Rules:
- Preserve all proper nouns, scientific/Latin terms, abbreviations, and URLs unchanged
- Preserve special characters like arrows and symbols
- Maintain original formatting, spacing, and punctuation style
- Return ONLY a valid JSON array of translated strings, in the same order as input
- The array must have exactly {len(texts)} elements

Input:
{json.dumps(texts, ensure_ascii=False)}"""

    try:
        response = await openai_client.chat.completions.create(
            model="o4-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        content = response.choices[0].message.content.strip()
        if content.startswith("```"):
            lines = content.split("\n")
            content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
            content = content.strip()
        result = json.loads(content)
        if isinstance(result, list) and len(result) == len(texts):
            return result
        while len(result) < len(texts):
            result.append(texts[len(result)])
        return result[:len(texts)]
    except Exception as e:
        logger.error(f"Translation batch failed: {e}")
        return texts


async def run_translation(job_id: str, target_language: str, tone: str):
    try:
        job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
        if not job:
            return

        segments = job["segments"]
        translatable = [(i, s) for i, s in enumerate(segments) if s["translatable"]]

        if not translatable:
            await db.translation_jobs.update_one(
                {"id": job_id},
                {"$set": {"status": "completed", "translated_segments": 0, "progress": 100}}
            )
            return

        batch_size = 20
        translated_count = 0

        for batch_start in range(0, len(translatable), batch_size):
            batch = translatable[batch_start:batch_start + batch_size]
            texts = [s["original"] for _, s in batch]
            translated_texts = await translate_batch(texts, target_language, tone)

            update_ops = {}
            for (orig_idx, _), translated_text in zip(batch, translated_texts):
                update_ops[f"segments.{orig_idx}.translated"] = translated_text

            translated_count += len(batch)
            progress = (translated_count / len(translatable)) * 100
            update_ops["translated_segments"] = translated_count
            update_ops["progress"] = progress

            await db.translation_jobs.update_one({"id": job_id}, {"$set": update_ops})
            await asyncio.sleep(0.5)

        # Non-translatable segments keep original
        final_ops = {}
        for i, s in enumerate(segments):
            if not s["translatable"]:
                final_ops[f"segments.{i}.translated"] = s["original"]

        final_ops["status"] = "completed"
        final_ops["progress"] = 100
        await db.translation_jobs.update_one({"id": job_id}, {"$set": final_ops})

        # Generate translated file and its preview images
        try:
            job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
            is_ocr = job.get("is_ocr", False)

            file_type = job["file_type"]
            original_path = job["file_path"]
            translated_path = str(UPLOAD_DIR / f"{job_id}_translated.{file_type}")

            if file_type == "pdf" and is_ocr:
                # For OCR PDFs, build translations dict with bbox, page, lines, and original text
                translations = {}
                for s in job["segments"]:
                    if s.get("translated") and s.get("ocr_bbox"):
                        translations[s["idx"]] = {
                            "text": s["translated"],
                            "bbox": s["ocr_bbox"],
                            "page": s["slide_num"] - 1,  # 0-indexed
                            "original": s["original"],
                            "lines": s.get("ocr_lines", []),
                        }
                rebuild_pdf(original_path, translated_path, translations, is_ocr=True)
            else:
                # Standard rebuild for text-based files
                translations = {}
                for s in job["segments"]:
                    if s.get("translated"):
                        translations[s["idx"]] = s["translated"]

                if file_type == "pptx":
                    rebuild_pptx(original_path, translated_path, translations)
                elif file_type == "docx":
                    rebuild_docx(original_path, translated_path, translations)
                elif file_type == "pdf":
                    rebuild_pdf(original_path, translated_path, translations, is_ocr=False)

            # Generate visual preview images for translated file
            translated_images = generate_slide_images(translated_path, file_type, job_id, "translated")
            translated_image_paths = [str(p) for p in translated_images]
            await db.translation_jobs.update_one(
                {"id": job_id},
                {"$set": {"translated_file_path": translated_path, "translated_images": translated_image_paths}}
            )
        except Exception as e:
            logger.error(f"Post-translation processing error: {e}")

        logger.info(f"Job {job_id} completed: {translated_count} segments translated")

    except Exception as e:
        logger.error(f"Translation job {job_id} failed: {e}")
        await db.translation_jobs.update_one(
            {"id": job_id},
            {"$set": {"status": "error", "error_message": str(e)}}
        )


# ════════════════════════════════════════
# API Routes
# ════════════════════════════════════════
@api_router.get("/")
async def root():
    return {"message": "PPTX Translator API"}


@api_router.post("/upload", response_model=JobInfo)
async def upload_file(file: UploadFile = File(...)):
    ext = Path(file.filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")

    job_id = str(uuid.uuid4())
    file_type = ext.lstrip(".")
    file_path = UPLOAD_DIR / f"{job_id}.{file_type}"

    content = await file.read()
    if len(content) > 100 * 1024 * 1024:
        raise HTTPException(400, "File too large (max 100MB)")

    with open(file_path, "wb") as f:
        f.write(content)

    try:
        is_ocr = False
        if file_type == "pptx":
            segments = extract_pptx_segments(str(file_path))
        elif file_type == "docx":
            segments = extract_docx_segments(str(file_path))
        elif file_type == "pdf":
            segments, is_ocr = extract_pdf_segments(str(file_path))
        else:
            raise HTTPException(400, "Unsupported file type")
    except HTTPException:
        raise
    except Exception as e:
        file_path.unlink(missing_ok=True)
        raise HTTPException(400, f"Failed to parse file: {str(e)}")

    # Generate original preview images in background
    original_images = generate_slide_images(str(file_path), file_type, job_id, "original")
    original_image_paths = [str(p) for p in original_images]

    job_doc = {
        "id": job_id,
        "filename": file.filename,
        "file_type": file_type,
        "file_path": str(file_path),
        "segments": segments,
        "total_segments": len(segments),
        "translated_segments": 0,
        "progress": 0,
        "status": "ready",
        "error_message": None,
        "original_images": original_image_paths,
        "translated_images": [],
        "is_ocr": is_ocr,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.translation_jobs.insert_one(job_doc)

    return JobInfo(id=job_id, filename=file.filename, file_type=file_type, total_segments=len(segments), status="ready")


@api_router.post("/translate/{job_id}")
async def start_translation(job_id: str, req: TranslateRequest):
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")
    if job["status"] == "translating":
        raise HTTPException(400, "Translation already in progress")

    await db.translation_jobs.update_one(
        {"id": job_id},
        {"$set": {
            "status": "translating",
            "progress": 0,
            "translated_segments": 0,
            "target_language": req.target_language,
            "tone": req.tone,
        }}
    )
    asyncio.create_task(run_translation(job_id, req.target_language, req.tone))
    return {"status": "translating", "job_id": job_id}


@api_router.get("/progress/{job_id}", response_model=ProgressResponse)
async def get_progress(job_id: str):
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")
    translatable_count = sum(1 for s in job["segments"] if s["translatable"])
    return ProgressResponse(
        id=job_id,
        status=job["status"],
        progress=job.get("progress", 0),
        total_segments=translatable_count,
        translated_segments=job.get("translated_segments", 0),
        error_message=job.get("error_message"),
    )


@api_router.get("/preview/{job_id}")
async def get_preview(job_id: str):
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")

    preview = []
    for s in job["segments"]:
        if s["translatable"] and s.get("translated"):
            preview.append({
                "idx": s["idx"],
                "slide_num": s["slide_num"],
                "original": s["original"],
                "translated": s["translated"],
            })
    return {"segments": preview, "status": job["status"]}


@api_router.get("/slides/{job_id}/{version}/{slide_index}")
async def get_slide_image(job_id: str, version: str, slide_index: int):
    """Serve slide preview images. version: 'original' or 'translated'."""
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")

    key = f"{version}_images"
    images = job.get(key, [])
    if slide_index < 0 or slide_index >= len(images):
        raise HTTPException(404, "Slide image not found")

    img_path = images[slide_index]
    if not Path(img_path).exists():
        raise HTTPException(404, "Image file not found")

    return FileResponse(img_path, media_type="image/png")


@api_router.get("/slides-info/{job_id}")
async def get_slides_info(job_id: str):
    """Return count of available slide images for original and translated."""
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")

    return {
        "original_count": len(job.get("original_images", [])),
        "translated_count": len(job.get("translated_images", [])),
        "file_type": job.get("file_type", "pptx"),
    }


@api_router.get("/download/{job_id}")
async def download_translated(job_id: str):
    job = await db.translation_jobs.find_one({"id": job_id}, {"_id": 0})
    if not job:
        raise HTTPException(404, "Job not found")
    if job["status"] != "completed":
        raise HTTPException(400, "Translation not yet completed")

    file_type = job["file_type"]
    translated_path = job.get("translated_file_path")

    if not translated_path or not Path(translated_path).exists():
        # Rebuild on the fly
        original_path = job["file_path"]
        if not Path(original_path).exists():
            raise HTTPException(404, "Original file not found")

        is_ocr = job.get("is_ocr", False)
        translated_path = str(UPLOAD_DIR / f"{job_id}_translated.{file_type}")

        if file_type == "pdf" and is_ocr:
            translations = {}
            for s in job["segments"]:
                if s.get("translated") and s.get("ocr_bbox"):
                    translations[s["idx"]] = {
                        "text": s["translated"],
                        "bbox": s["ocr_bbox"],
                        "page": s["slide_num"] - 1,
                        "original": s["original"],
                        "lines": s.get("ocr_lines", []),
                    }
            rebuild_pdf(original_path, translated_path, translations, is_ocr=True)
        else:
            translations = {}
            for s in job["segments"]:
                if s.get("translated"):
                    translations[s["idx"]] = s["translated"]

            if file_type == "pptx":
                rebuild_pptx(original_path, translated_path, translations)
            elif file_type == "docx":
                rebuild_docx(original_path, translated_path, translations)
            elif file_type == "pdf":
                rebuild_pdf(original_path, translated_path, translations, is_ocr=False)

    original_name = Path(job["filename"]).stem
    lang = job.get("target_language", "translated")
    download_name = f"{original_name}_{lang}.{file_type}"

    media_types = {
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pdf": "application/pdf",
    }

    return FileResponse(
        path=translated_path,
        filename=download_name,
        media_type=media_types.get(file_type, "application/octet-stream"),
    )


# ════════════════════════════════════════
# PDF Editor endpoints
# ════════════════════════════════════════

EDITOR_DIR = Path("/tmp/pdf_editor")
EDITOR_DIR.mkdir(exist_ok=True)


class EditorEditItem(BaseModel):
    type: str  # "text" or "whiteout"
    page: int
    x: float
    y: float
    width: float
    height: float
    text: Optional[str] = None
    fontSize: Optional[float] = 16
    fontColor: Optional[str] = "#000000"
    bold: Optional[bool] = False
    italic: Optional[bool] = False
    backgroundColor: Optional[str] = "#ffffff"


class EditorSaveRequest(BaseModel):
    edits: List[EditorEditItem]


@api_router.post("/editor/upload")
async def editor_upload(file: UploadFile = File(...)):
    ext = Path(file.filename).suffix.lower()
    if ext != ".pdf":
        raise HTTPException(status_code=400, detail="Only PDF files are supported in the editor")

    job_id = str(uuid.uuid4())
    job_dir = EDITOR_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    file_path = job_dir / file.filename
    content = await file.read()

    # File size limit (50MB)
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File size exceeds 50MB limit")

    with open(file_path, "wb") as f:
        f.write(content)

    # Open PDF and get page count
    doc = fitz.open(str(file_path))
    page_count = len(doc)

    # Render each page as PNG
    pages_dir = job_dir / "pages"
    pages_dir.mkdir(exist_ok=True)
    for i in range(page_count):
        page = doc[i]
        # Render at 2x for quality
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        pix.save(str(pages_dir / f"page_{i}.png"))

    # Store page dimensions for coordinate mapping
    page_dims = []
    for i in range(page_count):
        page = doc[i]
        page_dims.append({"width": page.rect.width, "height": page.rect.height})

    doc.close()

    # Save metadata
    meta = {
        "job_id": job_id,
        "filename": file.filename,
        "file_path": str(file_path),
        "page_count": page_count,
        "page_dims": page_dims,
    }
    with open(job_dir / "meta.json", "w") as f:
        json.dump(meta, f)

    return {
        "job_id": job_id,
        "filename": file.filename,
        "page_count": page_count,
        "page_dims": page_dims,
    }


@api_router.get("/editor/page/{job_id}/{page_num}")
async def editor_get_page(job_id: str, page_num: int):
    job_dir = EDITOR_DIR / job_id
    page_path = job_dir / "pages" / f"page_{page_num}.png"
    if not page_path.exists():
        raise HTTPException(status_code=404, detail="Page not found")
    return FileResponse(str(page_path), media_type="image/png")


@api_router.post("/editor/save/{job_id}")
async def editor_save(job_id: str, request: EditorSaveRequest):
    job_dir = EDITOR_DIR / job_id
    meta_path = job_dir / "meta.json"
    if not meta_path.exists():
        raise HTTPException(status_code=404, detail="Editor job not found")

    with open(meta_path) as f:
        meta = json.load(f)

    doc = fitz.open(meta["file_path"])

    for edit in request.edits:
        if edit.page < 0 or edit.page >= len(doc):
            continue
        page = doc[edit.page]

        if edit.type == "whiteout":
            # Draw white rectangle to cover content
            hex_color = edit.backgroundColor or "#ffffff"
            r = int(hex_color[1:3], 16) / 255
            g = int(hex_color[3:5], 16) / 255
            b = int(hex_color[5:7], 16) / 255
            rect = fitz.Rect(edit.x, edit.y, edit.x + edit.width, edit.y + edit.height)
            page.draw_rect(rect, color=None, fill=(r, g, b))

        elif edit.type == "text":
            # Add text at position
            hex_color = edit.fontColor or "#000000"
            r = int(hex_color[1:3], 16) / 255
            g = int(hex_color[3:5], 16) / 255
            b = int(hex_color[5:7], 16) / 255

            font_size = edit.fontSize or 16
            text = edit.text or ""

            # Use built-in font
            fontname = "helv"
            if edit.bold and edit.italic:
                fontname = "hebi"
            elif edit.bold:
                fontname = "hebo"
            elif edit.italic:
                fontname = "heit"

            # Insert text
            text_point = fitz.Point(edit.x, edit.y + font_size)
            page.insert_text(
                text_point,
                text,
                fontsize=font_size,
                fontname=fontname,
                color=(r, g, b),
            )

    # Save edited PDF
    output_path = job_dir / "edited.pdf"
    doc.save(str(output_path))
    doc.close()

    # Re-render pages for preview
    doc = fitz.open(str(output_path))
    pages_dir = job_dir / "pages_edited"
    pages_dir.mkdir(exist_ok=True)
    for i in range(len(doc)):
        page = doc[i]
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        pix.save(str(pages_dir / f"page_{i}.png"))
    doc.close()

    return {"status": "saved", "download_url": f"/api/editor/download/{job_id}"}


@api_router.get("/editor/page-edited/{job_id}/{page_num}")
async def editor_get_page_edited(job_id: str, page_num: int):
    job_dir = EDITOR_DIR / job_id
    page_path = job_dir / "pages_edited" / f"page_{page_num}.png"
    if not page_path.exists():
        raise HTTPException(status_code=404, detail="Edited page not found")
    return FileResponse(str(page_path), media_type="image/png")


@api_router.get("/editor/download/{job_id}")
async def editor_download(job_id: str):
    job_dir = EDITOR_DIR / job_id
    meta_path = job_dir / "meta.json"
    if not meta_path.exists():
        raise HTTPException(status_code=404, detail="Editor job not found")

    with open(meta_path) as f:
        meta = json.load(f)

    output_path = job_dir / "edited.pdf"
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="No edits saved yet")

    filename = meta["filename"].replace(".pdf", "_edited.pdf")
    return FileResponse(str(output_path), media_type="application/pdf", filename=filename)


app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
