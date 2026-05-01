#!/usr/bin/env python3
"""
Backend API Testing for PPTX Translator - Iteration 2
Tests all API endpoints with .pptx, .docx, .pdf support and visual preview
"""

import requests
import sys
import os
import time
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import fitz  # PyMuPDF for PDF creation

# Use the public backend URL from frontend .env
BACKEND_URL = "https://slide-translator-app.preview.emergentagent.com"
API_BASE = f"{BACKEND_URL}/api"

class PPTXTranslatorTester:
    def __init__(self):
        self.tests_run = 0
        self.tests_passed = 0
        self.job_ids = {}  # Store job IDs for different file types
        self.test_files = {}  # Store test file paths

    def log_test(self, name, success, details=""):
        """Log test result"""
        self.tests_run += 1
        if success:
            self.tests_passed += 1
            print(f"✅ {name} - PASSED {details}")
        else:
            print(f"❌ {name} - FAILED {details}")
        return success

    def create_test_files(self):
        """Create test files for all supported formats"""
        try:
            # Create PPTX test file
            prs = Presentation()
            slide1 = prs.slides.add_slide(prs.slide_layouts[1])
            slide1.shapes.title.text = "Test Presentation"
            content = slide1.placeholders[1].text_frame
            content.text = "This is a test slide for translation."
            
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Second Slide"
            content2 = slide2.placeholders[1].text_frame
            content2.text = "Another text segment to translate."
            
            self.test_files['pptx'] = "/tmp/test_presentation.pptx"
            prs.save(self.test_files['pptx'])
            
            # Create DOCX test file
            doc = Document()
            doc.add_heading('Test Document', 0)
            doc.add_paragraph('This is a test document for translation.')
            doc.add_heading('Second Section', level=1)
            doc.add_paragraph('Another paragraph to translate.')
            
            self.test_files['docx'] = "/tmp/test_document.docx"
            doc.save(self.test_files['docx'])
            
            # Create PDF test file
            pdf_doc = fitz.open()
            page1 = pdf_doc.new_page()
            page1.insert_text((72, 72), "Test PDF Document", fontsize=16)
            page1.insert_text((72, 120), "This is a test PDF for translation.", fontsize=12)
            
            page2 = pdf_doc.new_page()
            page2.insert_text((72, 72), "Second Page", fontsize=16)
            page2.insert_text((72, 120), "Another text block to translate.", fontsize=12)
            
            self.test_files['pdf'] = "/tmp/test_document.pdf"
            pdf_doc.save(self.test_files['pdf'])
            pdf_doc.close()
            
            return True
        except Exception as e:
            print(f"Failed to create test files: {e}")
            return False

    def test_health_check(self):
        """Test GET /api/ health check endpoint"""
        try:
            response = requests.get(f"{API_BASE}/", timeout=10)
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            if success:
                data = response.json()
                details += f", Response: {data}"
            return self.log_test("Health Check", success, details)
        except Exception as e:
            return self.log_test("Health Check", False, f"Error: {str(e)}")

    def test_upload_file(self, file_type):
        """Test POST /api/upload endpoint for different file types"""
        if file_type not in self.test_files or not os.path.exists(self.test_files[file_type]):
            return self.log_test(f"Upload {file_type.upper()}", False, "Test file not found")
        
        try:
            mime_types = {
                'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'pdf': 'application/pdf'
            }
            
            with open(self.test_files[file_type], 'rb') as f:
                files = {'file': (f'test.{file_type}', f, mime_types[file_type])}
                response = requests.post(f"{API_BASE}/upload", files=files, timeout=30)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                data = response.json()
                self.job_ids[file_type] = data.get('id')
                file_type_returned = data.get('file_type')
                details += f", Job ID: {self.job_ids[file_type]}, File Type: {file_type_returned}, Segments: {data.get('total_segments')}"
                # Verify file_type is correctly detected
                if file_type_returned != file_type:
                    success = False
                    details += f" (Expected {file_type}, got {file_type_returned})"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Upload {file_type.upper()}", success, details)
        except Exception as e:
            return self.log_test(f"Upload {file_type.upper()}", False, f"Error: {str(e)}")

    def test_upload_pptx(self):
        """Test POST /api/upload endpoint for PPTX"""
        return self.test_upload_file('pptx')
    
    def test_upload_docx(self):
        """Test POST /api/upload endpoint for DOCX"""
        return self.test_upload_file('docx')
    
    def test_upload_pdf(self):
        """Test POST /api/upload endpoint for PDF"""
        return self.test_upload_file('pdf')

    def test_upload_invalid_file(self):
        """Test upload with invalid file type"""
        try:
            # Create a fake txt file
            fake_file = b"This is not a PPTX file"
            files = {'file': ('test.txt', fake_file, 'text/plain')}
            response = requests.post(f"{API_BASE}/upload", files=files, timeout=10)
            
            # Should return 400 for invalid file type
            success = response.status_code == 400
            details = f"Status: {response.status_code}"
            return self.log_test("Upload Invalid File", success, details)
        except Exception as e:
            return self.log_test("Upload Invalid File", False, f"Error: {str(e)}")

    def test_start_translation(self, file_type='pptx'):
        """Test POST /api/translate/{job_id} endpoint"""
        if file_type not in self.job_ids:
            return self.log_test(f"Start Translation ({file_type.upper()})", False, "No job ID available")
        
        try:
            payload = {
                "target_language": "Spanish",
                "tone": "formal"
            }
            response = requests.post(f"{API_BASE}/translate/{self.job_ids[file_type]}", json=payload, timeout=10)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                data = response.json()
                details += f", Response: {data}"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Start Translation ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Start Translation ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_get_progress(self, file_type='pptx'):
        """Test GET /api/progress/{job_id} endpoint"""
        if file_type not in self.job_ids:
            return self.log_test(f"Get Progress ({file_type.upper()})", False, "No job ID available")
        
        try:
            response = requests.get(f"{API_BASE}/progress/{self.job_ids[file_type]}", timeout=10)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                data = response.json()
                details += f", Progress: {data.get('progress')}%, Status: {data.get('status')}"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Get Progress ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Get Progress ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_wait_for_completion(self, file_type='pptx', max_wait=120):
        """Wait for translation to complete and test progress polling"""
        if file_type not in self.job_ids:
            return self.log_test(f"Wait for Completion ({file_type.upper()})", False, "No job ID available")
        
        print(f"⏳ Waiting for {file_type.upper()} translation to complete (max {max_wait}s)...")
        start_time = time.time()
        
        try:
            while time.time() - start_time < max_wait:
                response = requests.get(f"{API_BASE}/progress/{self.job_ids[file_type]}", timeout=10)
                if response.status_code != 200:
                    return self.log_test(f"Wait for Completion ({file_type.upper()})", False, f"Progress check failed: {response.status_code}")
                
                data = response.json()
                status = data.get('status')
                progress = data.get('progress', 0)
                
                print(f"   Progress: {progress}%, Status: {status}")
                
                if status == "completed":
                    return self.log_test(f"Wait for Completion ({file_type.upper()})", True, f"Completed in {time.time() - start_time:.1f}s")
                elif status == "error":
                    error_msg = data.get('error_message', 'Unknown error')
                    return self.log_test(f"Wait for Completion ({file_type.upper()})", False, f"Translation failed: {error_msg}")
                
                time.sleep(3)  # Wait 3 seconds between checks
            
            return self.log_test(f"Wait for Completion ({file_type.upper()})", False, f"Timeout after {max_wait}s")
        except Exception as e:
            return self.log_test(f"Wait for Completion ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_get_preview(self, file_type='pptx'):
        """Test GET /api/preview/{job_id} endpoint"""
        if file_type not in self.job_ids:
            return self.log_test(f"Get Preview ({file_type.upper()})", False, "No job ID available")
        
        try:
            response = requests.get(f"{API_BASE}/preview/{self.job_ids[file_type]}", timeout=10)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                data = response.json()
                segments = data.get('segments', [])
                details += f", Segments: {len(segments)}"
                if segments:
                    details += f", Sample: '{segments[0].get('original', '')}' -> '{segments[0].get('translated', '')}'"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Get Preview ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Get Preview ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_slides_info(self, file_type='pptx'):
        """Test GET /api/slides-info/{job_id} endpoint - NEW in iteration 2"""
        if file_type not in self.job_ids:
            return self.log_test(f"Get Slides Info ({file_type.upper()})", False, "No job ID available")
        
        try:
            response = requests.get(f"{API_BASE}/slides-info/{self.job_ids[file_type]}", timeout=10)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                data = response.json()
                original_count = data.get('original_count', 0)
                translated_count = data.get('translated_count', 0)
                returned_file_type = data.get('file_type', '')
                details += f", Original: {original_count}, Translated: {translated_count}, File Type: {returned_file_type}"
                
                # Verify file type matches
                if returned_file_type != file_type:
                    success = False
                    details += f" (Expected {file_type}, got {returned_file_type})"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Get Slides Info ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Get Slides Info ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_slide_image(self, file_type='pptx'):
        """Test GET /api/slides/{job_id}/original/0 endpoint - NEW in iteration 2"""
        if file_type not in self.job_ids:
            return self.log_test(f"Get Slide Image ({file_type.upper()})", False, "No job ID available")
        
        try:
            # Test original slide image
            response = requests.get(f"{API_BASE}/slides/{self.job_ids[file_type]}/original/0", timeout=15)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                content_type = response.headers.get('content-type', '')
                content_length = len(response.content)
                details += f", Content-Type: {content_type}, Size: {content_length} bytes"
                
                # Verify it's a PNG image
                if 'image/png' in content_type and content_length > 1000:
                    details += " (Valid PNG)"
                else:
                    success = False
                    details += " (Invalid image)"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Get Slide Image ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Get Slide Image ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_download_translated(self, file_type='pptx'):
        """Test GET /api/download/{job_id} endpoint"""
        if file_type not in self.job_ids:
            return self.log_test(f"Download Translated ({file_type.upper()})", False, "No job ID available")
        
        try:
            response = requests.get(f"{API_BASE}/download/{self.job_ids[file_type]}", timeout=30)
            
            success = response.status_code == 200
            details = f"Status: {response.status_code}"
            
            if success:
                content_type = response.headers.get('content-type', '')
                content_length = len(response.content)
                details += f", Content-Type: {content_type}, Size: {content_length} bytes"
                
                # Verify content type matches file type
                expected_types = {
                    'pptx': 'presentation',
                    'docx': 'wordprocessing',
                    'pdf': 'pdf'
                }
                
                if expected_types[file_type] in content_type or content_length > 1000:
                    details += f" (Valid {file_type.upper()})"
                else:
                    success = False
                    details += f" (Invalid {file_type.upper()})"
            else:
                details += f", Error: {response.text}"
                
            return self.log_test(f"Download Translated ({file_type.upper()})", success, details)
        except Exception as e:
            return self.log_test(f"Download Translated ({file_type.upper()})", False, f"Error: {str(e)}")

    def test_invalid_job_id(self):
        """Test endpoints with invalid job ID"""
        fake_job_id = "invalid-job-id-12345"
        
        try:
            # Test progress with invalid ID
            response = requests.get(f"{API_BASE}/progress/{fake_job_id}", timeout=10)
            success1 = response.status_code == 404
            
            # Test preview with invalid ID
            response = requests.get(f"{API_BASE}/preview/{fake_job_id}", timeout=10)
            success2 = response.status_code == 404
            
            # Test download with invalid ID
            response = requests.get(f"{API_BASE}/download/{fake_job_id}", timeout=10)
            success3 = response.status_code == 404
            
            success = success1 and success2 and success3
            details = f"Progress: {success1}, Preview: {success2}, Download: {success3}"
            
            return self.log_test("Invalid Job ID Handling", success, details)
        except Exception as e:
            return self.log_test("Invalid Job ID Handling", False, f"Error: {str(e)}")

    def run_all_tests(self):
        """Run all backend tests for iteration 2"""
        print("🚀 Starting PPTX Translator Backend Tests - Iteration 2")
        print(f"📡 Testing API at: {API_BASE}")
        print("🆕 New features: .docx/.pdf support + visual preview")
        print("=" * 60)
        
        # Create test files
        if not self.create_test_files():
            print("❌ Failed to create test files. Exiting.")
            return False
        
        # Test basic functionality
        basic_tests = [
            self.test_health_check,
            self.test_upload_invalid_file,
        ]
        
        for test in basic_tests:
            test()
            print()
        
        # Test each file type
        file_types = ['pptx', 'docx', 'pdf']
        
        for file_type in file_types:
            print(f"🔄 Testing {file_type.upper()} workflow...")
            
            # Upload and start translation
            if not self.test_upload_file(file_type):
                continue
            if not self.test_start_translation(file_type):
                continue
            if not self.test_get_progress(file_type):
                continue
            if not self.test_wait_for_completion(file_type):
                continue
                
            # Test new visual preview endpoints
            self.test_slides_info(file_type)
            self.test_slide_image(file_type)
            
            # Test existing endpoints
            self.test_get_preview(file_type)
            self.test_download_translated(file_type)
            
            print()
        
        # Test error handling
        self.test_invalid_job_id()
        
        # Cleanup
        for file_path in self.test_files.values():
            if os.path.exists(file_path):
                os.remove(file_path)
        
        # Summary
        print("=" * 60)
        print(f"📊 Backend Tests Summary: {self.tests_passed}/{self.tests_run} passed")
        success_rate = (self.tests_passed / self.tests_run) * 100 if self.tests_run > 0 else 0
        print(f"📈 Success Rate: {success_rate:.1f}%")
        
        return self.tests_passed == self.tests_run

def main():
    """Main test runner"""
    tester = PPTXTranslatorTester()
    success = tester.run_all_tests()
    return 0 if success else 1

if __name__ == "__main__":
    sys.exit(main())